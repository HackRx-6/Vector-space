import os
import io
import zipfile
from typing import Optional, Tuple, IO, Any
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed

# PDF Processing
import fitz  # PyMuPDF

# DOCX Processing
import docx

# Excel Processing
import pandas as pd

# PPTX Processing
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation

# BIN Processing
import numpy as np

# Email (.eml)
from email import message_from_bytes
from bs4 import BeautifulSoup

# Load environment variables
from dotenv import load_dotenv
load_dotenv()


# ------------------- Basic Utilities -------------------
def basic_clean(text: str) -> str:
    return " ".join(text.split())


# ------------------- PDF Processing -------------------
def _extract_text_from_pdf_batch(pdf_bytes: bytes, start_page: int, end_page: int) -> Tuple[int, str]:
    batch_text = ""
    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            for i in range(start_page, end_page):
                page = doc.load_page(i)
                text = page.get_text() or ""
                if text:
                    batch_text += f"--- Page {i + 1} ---\n{text}\n\n"
            return start_page, batch_text
    except Exception as e:
        print(f"Error processing PDF batch starting at page {start_page}: {e}")
        return start_page, ""


def _process_pdf(pdf_bytes: bytes, batch_size: int) -> Optional[str]:
    try:
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            num_pages = doc.page_count
            print("PDF detected.")
            
            # Check for images
            has_images = any(len(page.get_images(full=True)) > 0 for page in doc)
            if has_images:
                print("PDF contains images.")

        if num_pages == 0:
            return ""

        tasks = []
        for i in range(0, num_pages, batch_size):
            start = i
            end = min(i + batch_size, num_pages)
            tasks.append((pdf_bytes, start, end))

        batch_results = {}
        with ProcessPoolExecutor() as executor:
            futures = {executor.submit(_extract_text_from_pdf_batch, *task): task[1] for task in tasks}
            for future in as_completed(futures):
                start_page = futures[future]
                try:
                    _, text = future.result()
                    batch_results[start_page] = text
                except Exception as exc:
                    print(f'PDF Batch starting at {start_page} generated an exception: {exc}')

        full_text = "".join(batch_results[start_page] for start_page in sorted(batch_results.keys()))
        return basic_clean(full_text)
    except Exception as e:
        print(f"An error occurred while processing the PDF: {e}")
        return None


# ------------------- DOCX Processing -------------------
def _process_docx(docx_bytes: bytes) -> Optional[str]:
    try:
        doc = docx.Document(io.BytesIO(docx_bytes))
        print("DOCX detected.")

        # Check for images
        if len(doc.inline_shapes) > 0:
            print("DOCX contains images.")

        full_text = "\n".join([para.text for para in doc.paragraphs])
        return basic_clean(full_text)
    except Exception as e:
        print(f"Error processing DOCX: {e}")
        return None


# ------------------- EML Processing -------------------
def _process_eml(eml_bytes: bytes) -> Optional[str]:
    try:
        msg = message_from_bytes(eml_bytes)
        body = ""

        def extract_text(part):
            content_type = part.get_content_type()
            charset = part.get_content_charset() or 'utf-8'
            if content_type == "text/plain":
                return part.get_payload(decode=True).decode(charset, errors='replace')
            elif content_type == "text/html":
                html = part.get_payload(decode=True).decode(charset, errors='replace')
                return BeautifulSoup(html, 'html.parser').get_text()
            return ""

        if msg.is_multipart():
            for part in msg.walk():
                body += extract_text(part)
        else:
            body = extract_text(msg)

        return body.strip()
    except Exception as e:
        print(f"Error processing EML: {e}")
        return None


# ------------------- Excel Processing -------------------
def process_excel(file_like: IO[bytes]) -> str:
    try:
        df = pd.read_excel(file_like, engine="openpyxl")

        if df.empty:
            return "Excel file is empty."

        # Convert to JSON (records format)
        json_output = df.to_json(orient="records")
        return json_output

    except Exception as e:
        return f"Error processing Excel file: {e}"


# ------------------- PPTX Processing -------------------
def extract_text_from_pptx_slide_fast(slide_data: Tuple[int, Any]) -> Tuple[int, str]:
    slide_num, slide = slide_data
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            texts.append(shape.text.strip())
        elif hasattr(shape, "table"):
            for row in shape.table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    texts.append(row_text)
    return slide_num, f"\n--- Slide {slide_num + 1} ---\n" + '\n'.join(texts) if texts else ""


def process_pptx(pptx_bytes: bytes) -> Optional[str]:
    try:
        pptx_stream = io.BytesIO(pptx_bytes)
        presentation = Presentation(pptx_stream)

        # Check for images
        has_images = any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for slide in presentation.slides for shape in slide.shapes)
        if has_images:
            print("PPTX contains images.")

        slide_tasks = [(i, slide) for i, slide in enumerate(presentation.slides)]
        slide_results = {}

        with ThreadPoolExecutor(max_workers=os.cpu_count() * 4) as executor:
            futures = {executor.submit(extract_text_from_pptx_slide_fast, task): task[0] for task in slide_tasks}
            for future in as_completed(futures):
                slide_num = futures[future]
                try:
                    _, text = future.result()
                    if text.strip():
                        slide_results[slide_num] = text
                except Exception as e:
                    print(f"Error processing slide {slide_num + 1}: {e}")

        full_text = ''.join(slide_results.get(i, '') for i in range(len(presentation.slides)))
        final = ' '.join(full_text.split()) if full_text else "No text found in PPTX"
        final = basic_clean(final)
        return final
    
    except Exception as e:
        print(f"Error processing PPTX: {e}")
        return None


# ------------------- BIN Processing -------------------
def process_bin(bin_bytes: bytes) -> Optional[str]:
    try:
        data_array = np.frombuffer(bin_bytes, dtype=np.uint8)
        printable_mask = ((data_array >= 32) & (data_array <= 126)) | np.isin(data_array, [9, 10, 13])
        ascii_chars = data_array[printable_mask]

        if len(ascii_chars) > len(bin_bytes) * 0.3:
            readable_text = ''.join(chr(c) for c in ascii_chars)
            if len(readable_text) > 20:
                result = ' '.join(readable_text.split())
                result.encode('utf-8')
                return result

        encodings = ['utf-8', 'latin-1', 'ascii', 'cp1252', 'iso-8859-1']

        def try_decode_fast(encoding):
            try:
                decoded = bin_bytes.decode(encoding, errors='ignore')
                readable = ''.join(c for c in decoded if c.isprintable() or c.isspace() or ord(c) > 127)
                return encoding, readable if len(readable) > 20 else ""
            except:
                return encoding, ""

        with ThreadPoolExecutor(max_workers=len(encodings)) as executor:
            results = list(executor.map(try_decode_fast, encodings))

        best_encoding, best_text = max(results, key=lambda x: len(x[1]))
        if best_text:
            result = ' '.join(best_text.split())
            result.encode('utf-8')
            return result
        else:
            return "Binary file - no readable text content"
    except:
        return None


# ------------------- ZIP Processing -------------------
def process_single_zip_file_sync(file_info: Tuple[str, bytes]) -> Tuple[str, str]:
    filename, file_bytes = file_info
    file_ext = os.path.splitext(filename)[1].lower()
    print(f"Processing file type: {file_ext}")

    try:
        if file_ext == '.pdf':
            return filename, _process_pdf(file_bytes, 25)
        elif file_ext == '.docx':
            return filename, _process_docx(file_bytes)
        elif file_ext == '.eml':
            return filename, _process_eml(file_bytes)
        elif file_ext in ['.xlsx', '.xls', '.xlsb']:
            file_like = io.BytesIO(file_bytes)
            return filename, process_excel(file_like)
        elif file_ext == '.bin':
            return filename, process_bin(file_bytes)
        elif file_ext == '.pptx':
            return filename, process_pptx(file_bytes)
        else:
            return filename, ""
    except Exception as e:
        print(f"Error processing {filename}: {e}")  # Add logging
        return filename, f"Error processing {filename}: {e}"


def process_zip_simple(zip_bytes: bytes) -> Optional[str]:
    try:
        # Use io.BytesIO which is a fully featured in-memory file-like object
        zip_file_like = io.BytesIO(zip_bytes)

        # The 'with' statement now correctly manages the ZipFile object
        with zipfile.ZipFile(zip_file_like, 'r') as zip_file:
            supported_extensions = {'.pdf', '.docx', '.eml', '.xlsx', '.bin', '.pptx', '.xls', '.xlsb'}
            combined_text = []
            
            for file_info in zip_file.infolist():
                if file_info.is_dir():
                    continue

                filename = file_info.filename
                file_ext = os.path.splitext(filename)[1].lower()

                if file_ext not in supported_extensions:
                    continue
                
                try:
                    file_bytes = zip_file.read(filename)
                    if len(file_bytes) == 0:
                        continue
                    
                    # Process file directly
                    result = process_single_zip_file_sync((filename, file_bytes))
                    if result and result[1] and not result[1].startswith("Error"):
                        combined_text.append(f"\n=== FILE: {result[0]} ===\n{result[1]}")
                
                except Exception as e:
                    print(f"Error processing {filename}: {e}")
                    continue

            return "\n\n".join(combined_text) if combined_text else "No text extracted"
    except Exception as e:
        # Catches errors like 'BadZipFile' if the bytes are not a valid zip
        print(f"An error occurred while processing the ZIP file: {e}")
        return f"Error processing ZIP file: {e}"